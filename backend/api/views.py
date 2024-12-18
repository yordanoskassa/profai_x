from django.shortcuts import render
from django.contrib.auth.models import User
from rest_framework import generics
from .serializers import UserSerializer, APIKeySerializer
from rest_framework.permissions import IsAuthenticated, AllowAny
from django.views.decorators.csrf import csrf_exempt
from .models import Key
import requests
from openai import OpenAI
from django.http import JsonResponse
from django.conf import settings
from rest_framework.views import APIView
from rest_framework.decorators import api_view, permission_classes
from rest_framework.permissions import AllowAny
from django.views.decorators.csrf import csrf_exempt
from django.views.decorators.http import require_POST
import hashlib
from django.contrib.auth.decorators import login_required
from rest_framework.parsers import JSONParser
from django.middleware.csrf import get_token
from rest_framework.exceptions import ParseError
import json
from pptx import Presentation
from rest_framework.response import Response
import openai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import requests
import json
from pptx.util import Inches
import io
import os
from django.http import FileResponse
from pptx.enum.text import PP_ALIGN



def get_csrf_token(request):
    token = get_token(request)
    return JsonResponse({'csrfToken': token})

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def save_api(request):
    if request.method == 'POST':
        try:
            data = JSONParser().parse(request)
            print("Request data received on backend:", data)
        except ParseError:
            return JsonResponse({'error': 'Invalid JSON data.'}, status=400)

        serializer = APIKeySerializer(data=data)
        if serializer.is_valid():
            if request.user.is_authenticated:
                api_key = serializer.save(author=request.user)
                return JsonResponse(
                    {
                        'message': 'API Key created successfully.',
                        'api_key': api_key.key
                    },
                    status=201
                )
            else:
                return JsonResponse({'error': 'User not authenticated.'}, status=401)
        return JsonResponse({'error': 'Invalid data.', 'details': serializer.errors}, status=400)
    return JsonResponse({'error': 'Invalid HTTP method.'}, status=405)

class APIKeyDelete(generics.DestroyAPIView):
    serializer_class = APIKeySerializer
    permission_classes = [IsAuthenticated]  # Only authenticated users

    def get_queryset(self):
        return Key.objects.filter(author=self.request.user)


class CreateUserView(generics.CreateAPIView):
    queryset = User.objects.all()
    serializer_class = UserSerializer
    permission_classes = [AllowAny]  # Allow anyone to register


    """ 
    try:
        user_key = Key.objects.get(author=request.user)
        api_key = user_key.key
    except Key.DoesNotExist:
        return JsonResponse({'error': 'API key not found for the user.'}, status=404) """


@api_view(['GET'])
@permission_classes([IsAuthenticated])  # Enforce authentication
def get_avatars(request):
    print("get avatars starting in views")

    api_key = settings.HEYGEN_API_KEY
    url = "https://api.heygen.com/v2/avatars"
    headers = {"accept": "application/json", "x-api-key": api_key}

    try:
        response = requests.get(url, headers=headers)
        response.raise_for_status()
        return JsonResponse(response.json(), status=200)
    except requests.exceptions.RequestException as e:
        return JsonResponse({'error': str(e)}, status=500)
    
    

@api_view(['GET'])
@permission_classes([AllowAny])  # Allow access without authentication
def get_voices(request):
    url = "https://api.heygen.com/v2/voices"
    api_key = settings.HEYGEN_API_KEY
    headers = {"accept": "application/json", "x-api-key": api_key}

    try:
        response = requests.get(url, headers=headers)
        response.raise_for_status()
        return JsonResponse(response.json(), status=200)
    except requests.exceptions.RequestException as e:
        return JsonResponse({'error': str(e)}, status=500)


@api_view(['POST'])
@permission_classes([IsAuthenticated])  # Enforce authentication
def generate_script(request):
    client = OpenAI(api_key=settings.CHATGPT_API_KEY)
    data = request.data
    avatar_name = data.get("selectedAvatar")
    content_prompt = data.get("contentPrompt") or "Provide a default prompt here."

    pre_prompt = (
        "You are a university-level instructional designer. Use the following content " + 
        "prompt to create an educational video script that adheres to sound instructional " +
        "design principles, and only includes text only. Please just have raw text, no headers of text at all. " + 
        "Do not have any indicators on who is speaking (no [Instructor]: label anywhere please, just the raw text).  "
    )
    messages = [
        {"role": "system", "content": f"Generating a script with text only."},
        {"role": "user", "content": pre_prompt + content_prompt},
    ]

    try:
        completion = client.chat.completions.create(model="gpt-3.5-turbo", messages=messages)
        generated_script = completion.choices[0].message.content
        return JsonResponse(
            {
                "message": "Script generated successfully",
                "avatar_name": avatar_name,
                "generated_script": generated_script,
            },
            status=200,
        )
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)



@api_view(['POST'])
@permission_classes([IsAuthenticated])  # Enforce authentication
def generate_slide_structure(request):
    """
    Generate a JSON structure for slides using OpenAI's GPT API based on a script.
    """
    print(" starting slide structure")
    client = OpenAI(api_key=settings.CHATGPT_API_KEY)
    data = request.data
    script = data.get("contentPrompt")
    try:
        pre_prompt = (
            "You are a university-level instructional designer tasked with creating a JSON representation of slides "
            "for a PowerPoint presentation based on the provided script. The output must follow these requirements:\n"
            "1. The JSON must strictly adhere to the following format:\n\n"
            "{\n"
            '  "slides": [\n'
            '    {\n'
            '      "title": "PowerPoint title",\n'
            '      "content": ["This presentation was generated by AI."]\n'
            '    },\n'
            '    {\n'
            '      "title": "Slide Title",\n'
            '      "content": ["Bullet point 1", "Bullet point 2", "Bullet point 3"]\n'
            '    },\n'
            '    {\n'
            '      "title": "Another Slide Title",\n'
            '      "content": ["Bullet point A", "Bullet point B"]\n'
            '    }\n'
            '  ]\n'
            "}\n\n"
            "2. The presentation must contain at least three slides:\n"
            "   - The first slide must always be a title slide indicating the presentation was generated by AI.\n"
            "   - The subsequent slides should be based on the script and include a title and up to 4 bullet points.\n"
            "   - If the script contains insufficient content for multiple slides, create additional slides using placeholders.\n"
            "3. Ensure all titles are concise, and the bullet points are clear and concise summaries of the script content.\n"
            "4. Ensure the JSON structure is valid and error-free.\n\n"
            "Script content follows below:\n\n"
        )

        messages = [
            {"role": "system", "content": "You are a JSON generator for PowerPoint slides with a rigorous format."},
            {"role": "user", "content": pre_prompt + script},
        ]


        # Retry logic
        max_attempts = 5
        for attempt in range(max_attempts):
            try:
                print("attempting chatgpt submission")
                completion = client.chat.completions.create(model="gpt-3.5-turbo", messages=messages)
                generated_script_outline = completion.choices[0].message.content
                slides_json = json.loads(generated_script_outline)

                if "slides" in slides_json and len(slides_json["slides"]) > 1:
                    #print(slides_json)
                    return JsonResponse(
                        {
                            "message": "Script outline generated successfully",
                            "generated_script_outline": generated_script_outline,
                        },
                        status=200,
                    )
                print("Response contained only one slide. Re-prompting ChatGPT...")

            except (json.JSONDecodeError, KeyError) as e:
                print(f"Error parsing response (Attempt {attempt + 1}): {e}")
                continue  # Retry
            

            except Exception as e:
                return JsonResponse({"error": str(e)}, status=500)
    
        raise Exception("Failed to generate valid slides after multiple attempts.")

    except Exception as e:
        # Return error response
        return Response({"error": str(e)}, status=500)



@api_view(['POST'])
@permission_classes([IsAuthenticated])  # Enforce authentication
def get_video_link(request):
    print("getting video id")
    video_id = request.data.get('video_id')
    if not video_id:
        print("no video id")
        return JsonResponse({'error': 'video_id is required'}, status=400)

    # Simulate fetching the video status and URL using the video_id
    # Replace this with actual logic to get the video status and URL
    url = f"https://api.heygen.com/v1/video_status.get?video_id={video_id}"

    api_key = settings.HEYGEN_API_KEY
    headers = {
        "accept": "application/json",
        "x-api-key": api_key
    }

    try:
        print(f"sending request to heygen with URL: {url}")  # Log the full URL
        response = requests.get(url, headers=headers)
        response.raise_for_status()
        print(f"Received response from heygen: {response.status_code}, {response.text}")  # Log the full response

        response_json = response.json()
        print(f"Response JSON: {response_json}")  # Log the JSON response

        status = response_json.get("data", {}).get("status")
        video_url = response_json.get("data", {}).get("video_url")

        if status == "completed":
            print(f"Video completed, URL: {video_url}")
            return JsonResponse({'video_url': video_url}, status=200)
        else:
            print(f"Video not ready, status: {status}")
            return JsonResponse({'status': status, 'message': 'Video is not ready yet'}, status=202)

    except requests.exceptions.RequestException as e:
        print(f"Request error: {str(e)}")  # Log the exception error
        return JsonResponse({'error': str(e)}, status=500)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_user_profile(request):
    user = request.user  # The currently authenticated user
    user_data = {
        'username': user.username,  # Return only the username (no password)
    }
    return Response(user_data)


'''
@api_view(['POST'])
@permission_classes([IsAuthenticated])  # Enforce authentication
def generate_video(request):
    print("Simulating video generation")

    # Simulated response
    simulated_response = {
        "error": None,
        "data": {
            "video_id": "a0becfdd531e4730ad08bd2ff36f8d69"
        }
    }

    print("Returning simulated response:", simulated_response)
    return JsonResponse(simulated_response, status=200)

'''

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def approve_slides(request):
    try:
        json_data = request.data  # Use request.data to get the JSON data
        prs = Presentation()
        # Set the presentation to widescreen (16:9) aspect ratio
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        print("starting python slide creation")
        
        # Correctly access the slides array from the nested dictionary
        slides = json_data["slides"]["slides"]["slides"]
        for index, slide_data in enumerate(slides):
            print(f"Processing slide {index + 1}")
            slide_title = slide_data["title"] if "title" in slide_data else "Untitled Slide"
            bullet_points = slide_data["content"] if "content" in slide_data else []
            # Check if it's the first slide (title slide)
            is_title_slide = index == 0
            create_slide(prs, slide_title, bullet_points, is_title_slide)

        # Save the presentation to an in-memory file
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        # Send the presentation as a file response
        response = FileResponse(output, as_attachment=True, filename='presentation.pptx')
        response['Content-Type'] = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        return response
    except Exception as e:
        print("Error:", str(e))  # Log the error
        return Response({"error": str(e)}, status=400)

def create_slide(prs, slide_title, bullet_points, is_title_slide):
    logo_path = os.path.join(os.path.dirname(__file__), 'gvsu_logo_slides.png')  # Path to the logo file

    if is_title_slide:
        # For title slide, use a layout with title and subtitle
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide Layout (index 0)
        title_shape = slide.shapes.title
        title_shape.text = slide_title
        # Add larger text or different formatting as needed (optional)
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)  # Larger title font size

        # Set subtitle to "This presentation was generated by AI"
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = "This presentation was generated by AI"
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(26)  # Set font size for subtitle

    else:
        # For regular slides, use a layout with title and content
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content Layout (index 1)
        
        # Adjust the title text box position and size to make it wider and set height
        title_shape = slide.shapes.title
        title_shape.text = slide_title
        title_shape.width = Inches(12)  # Make the title text box wider
        title_shape.height = Inches(1.5)  # Set height for the title text box

        # Move the title text box up higher and to the left a little bit
        title_shape.left = Inches(1.25)
        title_shape.top = Inches(0.5)

        # Align text to the left
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT

        # Add bullet points
        text_box = slide.placeholders[1]
        text_box.text = ""  # Clear any placeholder text

        for bullet_point in bullet_points:
            paragraph = text_box.text_frame.add_paragraph()
            paragraph.text = bullet_point.strip()
            paragraph.font.size = Pt(26)  # Set font size to 26 points

        # Adjust the size of the text box and position it correctly
        text_box.left = Inches(1)
        text_box.top = Inches(1.5)
        text_box.width = Inches(11)
        text_box.height = Inches(4.5)

    # Add the blue box at the bottom of the slide
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    box_height = Inches(1)  # Set the height of the blue box
    box_top = slide_height - box_height
    blue_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=0,
        top=box_top,
        width=slide_width,
        height=box_height,
    )
    blue_box.fill.solid()
    blue_box.fill.fore_color.rgb = RGBColor(0, 0, 255)
    blue_box.line.color.rgb = RGBColor(0, 0, 255)  # No border

    # Add logo to the bottom right corner of every slide

    logo_width = Inches(2.5)  # Adjust as needed for logo size and position
    logo_height = Inches(0.72)  # Adjust as needed for logo size and position
    logo_left = slide_width - logo_width - Inches(0.5)  # Adjust as needed for logo size and position
    logo_top = slide_height - logo_height - Inches(0.2)  # Adjust as needed for logo size and position



    slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width, height=logo_height)

    return slide

@api_view(['POST'])
@permission_classes([IsAuthenticated])  # Enforce authentication
def generate_video(request):
    print("starting video generation")
    url = "https://api.heygen.com/v2/video/generate"
    api_key = settings.HEYGEN_API_KEY
    script = request.data.get("script")
    title = request.data.get("title")
    avatar = request.data.get("selectedAvatar")
    voice = request.data.get("selectedVoice")


    headers = {
        "accept": "application/json",
        "content-type": "application/json",
        "x-api-key": api_key
    }

    # Extract data from the request
    script = request.data.get("script")
    title = request.data.get("title")
    avatar = request.data.get("selectedAvatar")
    voice = request.data.get("selectedVoice")

    print("Avatar Voice ID: " + str(voice))
    print("Avatar ID: " + str(avatar))

    payload = {
        "caption": False,
        "dimension": {"width": 1280, "height": 720},
        "title": title,
        "video_inputs": [
            {
                "character": {
                    "type": "avatar",
                    "avatar_id": avatar,
                    "scale": 1.0
                },
                "voice": {
                    "type": "text",
                    "voice_id": voice,  # Replace with actual voice ID
                    "input_text": script
                }
            }
        ]
    }

    headers = {
        "accept": "application/json",
        "content-type": "application/json",
        "x-api-key": api_key
    }

    response = requests.post(url, json=payload, headers=headers)
    print(response.text)

    try:
        print("parsing responce")
        response = requests.post(url, json=payload, headers=headers)
        response.raise_for_status()

        response_json = response.json()
        return JsonResponse(response_json, status=response.status_code)
    except requests.exceptions.RequestException as e:
        return JsonResponse({'error': str(e)}, status=500)






#### For when we have a real website domain, we can use webhooks to notify our 
# browser when a video is done generating as opposed to having a 
# user click a button over and over. 

#@csrf_exempt
#@require_POST
#def heygen_webhook(request):
#    data = request.json()
#    video_id = data.get('video_id')
#    status = data.get('status')
#    video_url = data.get('video_url')
#
#    if status == 'completed':
#        # Handle the completed video, e.g., save the video URL to the database
#        print(f"Video {video_id} is ready at {video_url}")
#        # You can also update the frontend via WebSocket or other means
#
#    return JsonResponse({'message': 'Webhook received'}, status=200)


#@api_view(['POST'])
#@permission_classes([IsAuthenticated])  # Enforce authentication
#def register_webhook(request):
#    url = "https://api.heygen.com/v2/webhooks"
#    api_key = settings.HEYGEN_API_KEY
#    headers = {
#        "accept": "application/json",
#        "content-type": "application/json",
#        "x-api-key": api_key
#    }
#    payload = {
#        "url": "http://127.0.0.1:8000/api/heygen_webhook/",
#        "event": "video.completed"
#    }
#
#    try:
#        response = requests.post(url, json=payload, headers=headers)
#        response.raise_for_status()
#        return JsonResponse({'message': 'Webhook registered successfully'}, status=200)
#    except requests.exceptions.RequestException as e:
#        return JsonResponse({'error': str(e)}, status=500)